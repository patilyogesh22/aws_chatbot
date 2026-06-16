"""
ingestion.py
Extracts raw text from any uploaded file and writes
chunk records into PostgreSQL (raw layer) for dbt to transform.
"""

import os
import json
import uuid
import csv
import boto3
import tempfile

from pathlib import Path
from typing import List, Dict

import psycopg2
from psycopg2.extras import execute_values

from app.config import PG_DSN, CHUNK_SIZE, CHUNK_OVERLAP


# ─────────────────────────────────────────────────────────────
# AWS S3 CONFIG
# ─────────────────────────────────────────────────────────────

AWS_REGION = os.getenv("AWS_REGION", "eu-north-1")
S3_BUCKET = os.getenv("S3_BUCKET")

s3 = boto3.client("s3", region_name=AWS_REGION)


# ─────────────────────────────────────────────────────────────
# TEXT EXTRACTION
# ─────────────────────────────────────────────────────────────

def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext in (".csv", ".tsv"):
        return _extract_csv(file_path)
    elif ext in (".xlsx", ".xls"):
        return _extract_excel(file_path)
    elif ext == ".json":
        return _extract_json(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext in (".txt", ".md", ".rst", ".log"):
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    else:
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: str) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    result = "\n".join(texts)

    print("PDF TEXT LENGTH:", len(result))
    print("HAS NUL:", "\x00" in result)

    return result


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_csv(path: str) -> str:
    lines = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append(", ".join(row))
    return "\n".join(lines)


def _extract_excel(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)

    lines = []
    for sheet in wb.worksheets:
        lines.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            line = ", ".join(str(v) for v in row if v is not None)
            if line.strip():
                lines.append(line)

    return "\n".join(lines)


def _extract_json(path: str) -> str:
    with open(path, encoding="utf-8") as f:
        return json.dumps(json.load(f), indent=2)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    texts = []

    for i, slide in enumerate(prs.slides, 1):
        texts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())

    return "\n".join(texts)


# ─────────────────────────────────────────────────────────────
# CHUNKING
# ─────────────────────────────────────────────────────────────


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
        i += chunk_size - overlap

    return [c for c in chunks if c.strip()]


# ─────────────────────────────────────────────────────────────
# POSTGRES
# ─────────────────────────────────────────────────────────────

def get_conn():
    return psycopg2.connect(PG_DSN)


def init_postgres():
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE EXTENSION IF NOT EXISTS vector;

                CREATE TABLE IF NOT EXISTS raw_chunks (
                    chunk_id TEXT PRIMARY KEY,
                    file_name TEXT NOT NULL,
                    file_path TEXT,
                    file_hash TEXT,              -- ✅ ADD THIS
                    chunk_index INTEGER NOT NULL,
                    chunk_text TEXT NOT NULL,
                    char_count INTEGER,
                    word_count INTEGER,
                    file_size BIGINT DEFAULT 0,
                    ingested_at TIMESTAMPTZ DEFAULT NOW()
                );

                CREATE INDEX IF NOT EXISTS idx_raw_chunks_file_name
                ON raw_chunks(file_name);

                CREATE INDEX IF NOT EXISTS idx_raw_chunks_file_hash
                ON raw_chunks(file_hash);
            """)
        conn.commit()

    print("[ingestion] PostgreSQL ready.")

def ingest_file(file_path: str, file_name: str = None, file_hash: str = None) -> List[Dict]:
    init_postgres()

    if file_name is None:
        file_name = os.path.basename(file_path)

    file_size = os.path.getsize(file_path)

    raw_text = extract_text(file_path)
    # Remove PostgreSQL-breaking characters
    raw_text = raw_text.replace("\x00", "")

    chunks = chunk_text(raw_text)

    records = []
    rows = []

    for idx, chunk in enumerate(chunks):
        chunk = chunk.replace("\x00", "")
        rec = {
            "chunk_id": str(uuid.uuid4()),
            "file_name": file_name,
            "file_path": file_path,
            "file_hash": file_hash,   # ✅ ADD THIS
            "chunk_index": idx,
            "chunk_text": chunk,
            "char_count": len(chunk),
            "word_count": len(chunk.split()),
            "file_size": file_size,
        }

        records.append(rec)

        rows.append((
            rec["chunk_id"],
            rec["file_name"],
            rec["file_path"],
            rec["file_hash"],   # ✅ ADD THIS
            rec["chunk_index"],
            rec["chunk_text"],
            rec["char_count"],
            rec["word_count"],
            rec["file_size"],
        ))

    with get_conn() as conn:
        with conn.cursor() as cur:
            if file_hash:
                cur.execute("DELETE FROM raw_chunks WHERE file_hash = %s", (file_hash,))
            else:
                cur.execute("DELETE FROM raw_chunks WHERE file_name = %s", (file_name,))
            execute_values(cur, """
                INSERT INTO raw_chunks (
                    chunk_id, file_name, file_path, file_hash,
                    chunk_index, chunk_text,
                    char_count, word_count, file_size
                )
                VALUES %s
            """, rows)
        conn.commit()

    return records
# ─────────────────────────────────────────────────────────────
# S3 INGESTION (IMPORTANT FIX)
# ─────────────────────────────────────────────────────────────

def ingest_file_from_s3_key(bucket: str, s3_key: str, file_name: str,file_hash: str = None) -> List[Dict]:
    """
    Download file from S3 → temp file → reuse ingestion pipeline
    """

    if not bucket:
        raise ValueError("S3 bucket is not configured")

    tmp_path = None

    try:
        suffix = Path(file_name).suffix

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=suffix
        ) as tmp:
            tmp_path = tmp.name
            
        # download from S3
        s3.download_file(bucket, s3_key, tmp_path)

        # IMPORTANT FIX: pass correct file_name + file_path
        return ingest_file(
            file_path=tmp_path,
            file_name=file_name,
            file_hash=file_hash
        )

    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)
            
def delete_file_chunks(file_name: str):
    """Remove all raw_chunks rows for a given file."""
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("DELETE FROM raw_chunks WHERE file_name = %s", (file_name,))
        conn.commit()
    print(f"[ingestion] Deleted raw_chunks for '{file_name}'")


def list_ingested_files() -> List[str]:
    """Return distinct file names present in raw_chunks."""
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT DISTINCT file_name FROM raw_chunks ORDER BY file_name")
            return [row[0] for row in cur.fetchall()]

def list_ingested_files_meta() -> List[Dict]:
    """Return file list with chunk count, file size, and upload date."""
    with get_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT
                    file_name,
                    COUNT(*)          AS chunk_count,
                    MAX(file_size)    AS file_size,
                    MAX(ingested_at)  AS uploaded_at
                FROM raw_chunks
                GROUP BY file_name
                ORDER BY MAX(ingested_at) DESC
            """)
            rows = cur.fetchall()
    return [
        {
            "name":        r[0],
            "chunks":      r[1],
            "size":        r[2] or 0,
            "uploaded_at": r[3].isoformat() if r[3] else None,
        }
        for r in rows
    ]
